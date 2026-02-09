"""Doküman Yönetimi API Routes - RAG Sistemi
Departman bazlı doküman CRUD, genişletilmiş format desteği,
URL/link öğrenme ve YouTube video öğrenme özellikleri.
"""

from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form
from sqlalchemy.ext.asyncio import AsyncSession
from pydantic import BaseModel, Field
from typing import List, Optional
from datetime import datetime
import io
import json
import re
import structlog

from app.db.database import get_db
from app.db.models import User
from app.api.routes.auth import get_current_user

logger = structlog.get_logger()

# RAG modülünü güvenli şekilde import et
try:
    from app.rag.vector_store import (
        add_document,
        search_documents,
        get_stats,
        delete_document,
        clear_all_documents,
        list_documents,
    )
    RAG_AVAILABLE = True
except ImportError:
    RAG_AVAILABLE = False

# Web scraping kütüphaneleri
try:
    import httpx
    from bs4 import BeautifulSoup
    WEB_SCRAPING_AVAILABLE = True
except ImportError:
    WEB_SCRAPING_AVAILABLE = False

# YouTube transcript kütüphanesi
try:
    from youtube_transcript_api import YouTubeTranscriptApi
    YOUTUBE_AVAILABLE = True
except ImportError:
    YOUTUBE_AVAILABLE = False

router = APIRouter()

# ═══════════════════════════════════════════════════════════════
# DESTEKLENEN DOSYA FORMATLARI (Genişletilmiş v2)
# ═══════════════════════════════════════════════════════════════
SUPPORTED_FORMATS = {
    # ── Metin dosyaları ──
    '.txt': 'text',
    '.md': 'markdown',
    '.csv': 'csv',
    '.json': 'json',
    '.xml': 'xml',
    '.html': 'html',
    '.htm': 'html',
    '.rtf': 'rtf',
    '.rst': 'restructuredtext',
    '.tex': 'latex',
    '.ini': 'ini',
    '.cfg': 'config',
    '.env': 'env',
    '.toml': 'toml',
    '.properties': 'properties',
    # ── Office dosyaları ──
    '.pdf': 'pdf',
    '.docx': 'docx',
    '.doc': 'doc',
    '.xlsx': 'excel',
    '.xls': 'excel',
    '.pptx': 'powerpoint',
    '.ppt': 'powerpoint',
    '.odt': 'odt',
    '.ods': 'ods',
    '.odp': 'odp',
    '.epub': 'epub',
    # ── Kod dosyaları ──
    '.py': 'python',
    '.js': 'javascript',
    '.ts': 'typescript',
    '.jsx': 'jsx',
    '.tsx': 'tsx',
    '.java': 'java',
    '.cs': 'csharp',
    '.cpp': 'cpp',
    '.c': 'c',
    '.h': 'c_header',
    '.hpp': 'cpp_header',
    '.sql': 'sql',
    '.yaml': 'yaml',
    '.yml': 'yaml',
    '.go': 'go',
    '.rb': 'ruby',
    '.php': 'php',
    '.swift': 'swift',
    '.kt': 'kotlin',
    '.scala': 'scala',
    '.rs': 'rust',
    '.r': 'r_lang',
    '.R': 'r_lang',
    '.sh': 'shell',
    '.bat': 'batch',
    '.ps1': 'powershell',
    '.dockerfile': 'dockerfile',
    '.vue': 'vue',
    '.svelte': 'svelte',
    '.graphql': 'graphql',
    '.gql': 'graphql',
    '.proto': 'protobuf',
    # ── E-posta dosyaları ──
    '.eml': 'email',
    '.msg': 'email',
    # ── Görüntü dosyaları (OCR) ──
    '.png': 'image',
    '.jpg': 'image',
    '.jpeg': 'image',
    '.gif': 'image',
    '.bmp': 'image',
    '.tiff': 'image',
    '.tif': 'image',
    '.webp': 'image',
    # ── Log dosyaları ──
    '.log': 'log',
}


class DocumentAddRequest(BaseModel):
    content: str
    source: str
    doc_type: str = "text"
    department: str = "Genel"


class SearchRequest(BaseModel):
    query: str
    department: str = "Genel"
    n_results: int = 3


class TeachRequest(BaseModel):
    content: str
    department: str = "Genel"
    title: str = "Hızlı Bilgi Girişi"


class LearnFromUrlRequest(BaseModel):
    """URL'den öğrenme isteği"""
    url: str = Field(..., description="Öğrenilecek web sayfasının URL'si")
    department: str = "Genel"
    title: Optional[str] = Field(None, description="Opsiyonel başlık, belirtilmezse sayfa başlığı kullanılır")


class LearnFromVideoRequest(BaseModel):
    """YouTube/Video platformundan öğrenme isteği"""
    url: str = Field(..., description="YouTube veya video platformu URL'si")
    department: str = "Genel"
    title: Optional[str] = Field(None, description="Opsiyonel başlık")
    language: str = Field("tr", description="Tercih edilen altyazı dili (tr, en, vb.)")


class DocumentResponse(BaseModel):
    content: str
    source: str
    relevance: float


class DocumentListItem(BaseModel):
    source: str
    type: str
    department: str
    author: str
    created_at: str
    chunk_count: int


def parse_user_departments(department_str: str) -> List[str]:
    """Kullanıcı departmanlarını parse et"""
    if not department_str:
        return []
    try:
        parsed = json.loads(department_str)
        return parsed if isinstance(parsed, list) else [department_str]
    except json.JSONDecodeError:
        return [department_str]


def extract_text_from_file(filename: str, file_content: bytes) -> tuple:
    """Dosyadan metin çıkar - Genişletilmiş format desteği"""
    filename_lower = filename.lower()
    content = ""
    doc_type = "text"
    
    # Uzantıyı bul
    ext = None
    for e in SUPPORTED_FORMATS:
        if filename_lower.endswith(e):
            ext = e
            doc_type = SUPPORTED_FORMATS[e]
            break
    
    if ext is None:
        raise HTTPException(
            status_code=400,
            detail=f"Desteklenmeyen dosya formatı. Desteklenen formatlar: {', '.join(SUPPORTED_FORMATS.keys())}"
        )
    
    try:
        # ── Metin tabanlı dosyalar ──
        text_types = [
            'text', 'markdown', 'csv', 'json', 'xml', 'html', 'python',
            'javascript', 'typescript', 'jsx', 'tsx', 'java', 'csharp', 
            'cpp', 'c', 'c_header', 'cpp_header', 'sql', 'yaml', 'log',
            'restructuredtext', 'latex', 'ini', 'config', 'env', 'toml',
            'properties', 'go', 'ruby', 'php', 'swift', 'kotlin', 'scala',
            'rust', 'r_lang', 'shell', 'batch', 'powershell', 'dockerfile',
            'vue', 'svelte', 'graphql', 'protobuf'
        ]
        if doc_type in text_types:
            try:
                content = file_content.decode('utf-8')
            except UnicodeDecodeError:
                content = file_content.decode('latin-1', errors='ignore')
        
        # ── PDF ──
        elif doc_type == 'pdf':
            try:
                from PyPDF2 import PdfReader
                pdf = PdfReader(io.BytesIO(file_content))
                content = "\n".join([page.extract_text() or "" for page in pdf.pages])
            except ImportError:
                raise HTTPException(status_code=500, detail="PyPDF2 yüklü değil")
        
        # ── Word (DOCX) ──
        elif doc_type == 'docx':
            try:
                from docx import Document
                doc = Document(io.BytesIO(file_content))
                content = "\n".join([para.text for para in doc.paragraphs])
            except ImportError:
                raise HTTPException(status_code=500, detail="python-docx yüklü değil")
        
        # ── Excel ──
        elif doc_type == 'excel':
            try:
                import pandas as pd
                df = pd.read_excel(io.BytesIO(file_content))
                content = df.to_string()
            except ImportError:
                raise HTTPException(status_code=500, detail="pandas/openpyxl yüklü değil")
        
        # ── PowerPoint ──
        elif doc_type == 'powerpoint':
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_content))
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_parts.append(shape.text)
                content = "\n".join(text_parts)
            except ImportError:
                raise HTTPException(status_code=500, detail="python-pptx yüklü değil")
        
        # ── RTF ──
        elif doc_type == 'rtf':
            try:
                from striprtf.striprtf import rtf_to_text
                raw = file_content.decode('utf-8', errors='ignore')
                content = rtf_to_text(raw)
            except ImportError:
                # Fallback: basit regex ile RTF tag'larını temizle
                raw = file_content.decode('utf-8', errors='ignore')
                content = re.sub(r'[\\{}\[\]]', '', raw)
                content = re.sub(r'\\[a-z]+\d*\s?', '', content)
        
        # ── ODT (OpenDocument Text) ──
        elif doc_type in ('odt', 'ods', 'odp'):
            try:
                import zipfile
                with zipfile.ZipFile(io.BytesIO(file_content)) as z:
                    if 'content.xml' in z.namelist():
                        xml_content = z.read('content.xml').decode('utf-8')
                        # XML tag'larını temizle
                        content = re.sub(r'<[^>]+>', ' ', xml_content)
                        content = re.sub(r'\s+', ' ', content).strip()
            except Exception as e:
                raise HTTPException(status_code=500, detail=f"ODT okuma hatası: {str(e)}")
        
        # ── EPUB ──
        elif doc_type == 'epub':
            try:
                import zipfile
                with zipfile.ZipFile(io.BytesIO(file_content)) as z:
                    text_parts = []
                    for name in z.namelist():
                        if name.endswith(('.html', '.xhtml', '.htm')):
                            html_content = z.read(name).decode('utf-8', errors='ignore')
                            # HTML tag'larını temizle
                            clean = re.sub(r'<[^>]+>', ' ', html_content)
                            clean = re.sub(r'\s+', ' ', clean).strip()
                            if clean:
                                text_parts.append(clean)
                    content = "\n\n".join(text_parts)
            except Exception as e:
                raise HTTPException(status_code=500, detail=f"EPUB okuma hatası: {str(e)}")
        
        # ── E-posta (.eml) ──
        elif doc_type == 'email' and ext == '.eml':
            try:
                import email
                from email import policy
                msg = email.message_from_bytes(file_content, policy=policy.default)
                parts = []
                if msg['Subject']:
                    parts.append(f"Konu: {msg['Subject']}")
                if msg['From']:
                    parts.append(f"Gönderen: {msg['From']}")
                if msg['Date']:
                    parts.append(f"Tarih: {msg['Date']}")
                parts.append("---")
                body = msg.get_body(preferencelist=('plain', 'html'))
                if body:
                    body_content = body.get_content()
                    if body.get_content_type() == 'text/html':
                        body_content = re.sub(r'<[^>]+>', ' ', body_content)
                    parts.append(body_content)
                content = "\n".join(parts)
            except Exception as e:
                content = file_content.decode('utf-8', errors='ignore')
        
        # ── Görüntü dosyaları (OCR) ──
        elif doc_type == 'image':
            try:
                from PIL import Image
                import pytesseract
                img = Image.open(io.BytesIO(file_content))
                content = pytesseract.image_to_string(img, lang='tur+eng')
            except ImportError:
                # OCR mevcut değilse, görüntü metadatasını kullan
                try:
                    from PIL import Image
                    img = Image.open(io.BytesIO(file_content))
                    content = f"[Görüntü dosyası: {filename}, Boyut: {img.size[0]}x{img.size[1]}, Format: {img.format}]"
                    logger.warning("pytesseract_not_available", file=filename, msg="OCR için pytesseract yüklü değil, sadece metadata kaydedildi")
                except ImportError:
                    raise HTTPException(status_code=500, detail="Pillow yüklü değil, görüntü dosyaları işlenemiyor")
        
        return content, doc_type
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Dosya okuma hatası: {str(e)}")


@router.get("/status")
async def get_rag_status(current_user: User = Depends(get_current_user)):
    """RAG sistemi durumu"""
    if not RAG_AVAILABLE:
        return {
            "available": False,
            "message": "RAG modülü kurulu değil."
        }
    
    stats = get_stats()
    return {
        "available": True,
        "document_count": stats.get("document_count", 0),
        "collection": stats.get("collection"),
        "embedding_model": stats.get("embedding_model"),
    }


@router.get("/documents/list")
async def get_documents_list(
    current_user: User = Depends(get_current_user)
):
    """Kullanıcının departmanına göre doküman listesi"""
    if not RAG_AVAILABLE:
        raise HTTPException(status_code=503, detail="RAG sistemi kullanılamıyor")
    
    # Admin/Manager tüm dokümanları görür
    if current_user.role in ['admin', 'manager']:
        documents = list_documents(department=None)
    else:
        # User sadece kendi departmanlarını görür
        user_depts = parse_user_departments(current_user.department)
        all_docs = []
        for dept in user_depts:
            docs = list_documents(department=dept)
            all_docs.extend(docs)
        # Benzersiz kaynak isimlerine göre filtrele
        seen = set()
        documents = []
        for doc in all_docs:
            if doc['source'] not in seen:
                seen.add(doc['source'])
                documents.append(doc)
    
    return {"documents": documents, "count": len(documents)}


@router.post("/documents/add")
async def add_document_text(
    request: DocumentAddRequest,
    current_user: User = Depends(get_current_user)
):
    """Metin dokümanı ekle"""
    if not RAG_AVAILABLE:
        raise HTTPException(status_code=503, detail="RAG sistemi kullanılamıyor")
    
    # Departman yetki kontrolü
    if current_user.role == 'user':
        user_depts = parse_user_departments(current_user.department)
        if request.department not in user_depts:
            raise HTTPException(status_code=403, detail="Bu departmana doküman ekleme yetkiniz yok")
    
    success = add_document(
        content=request.content,
        source=request.source,
        doc_type=request.doc_type,
        metadata={
            "department": request.department,
            "author": current_user.email,
            "created_at": str(datetime.utcnow())
        }
    )
    
    if success:
        return {"message": f"Doküman eklendi: {request.source}", "success": True}
    else:
        raise HTTPException(status_code=500, detail="Doküman eklenemedi")


@router.post("/documents/upload")
async def upload_document(
    file: UploadFile = File(...),
    department: str = Form("Genel"),
    current_user: User = Depends(get_current_user)
):
    """Dosya yükle - Çoklu format desteği"""
    if not RAG_AVAILABLE:
        raise HTTPException(status_code=503, detail="RAG sistemi kullanılamıyor")
    
    # Departman yetki kontrolü
    if current_user.role == 'user':
        user_depts = parse_user_departments(current_user.department)
        if department not in user_depts:
            raise HTTPException(status_code=403, detail="Bu departmana dosya yükleme yetkiniz yok")
    
    try:
        file_content = await file.read()
        
        if not file_content:
            raise HTTPException(status_code=400, detail="Dosya boş")
        
        content, doc_type = extract_text_from_file(file.filename, file_content)
        
        if not content.strip():
            raise HTTPException(status_code=400, detail="Dosyadan içerik çıkarılamadı")
        
        # Dokümanı ekle
        success = add_document(
            content=content,
            source=file.filename,
            doc_type=doc_type,
            metadata={
                "department": department,
                "author": current_user.email,
                "created_at": str(datetime.utcnow())
            }
        )
        
        if success:
            return {
                "message": f"Dosya yüklendi: {file.filename}",
                "success": True,
                "type": doc_type,
                "chars": len(content)
            }
        else:
            raise HTTPException(status_code=500, detail="Doküman eklenemedi")
            
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Dosya işleme hatası: {str(e)}")


@router.post("/documents/upload-multiple")
async def upload_multiple_documents(
    files: List[UploadFile] = File(...),
    department: str = Form("Genel"),
    current_user: User = Depends(get_current_user)
):
    """Birden fazla dosya yükle (klasör yükleme için)"""
    if not RAG_AVAILABLE:
        raise HTTPException(status_code=503, detail="RAG sistemi kullanılamıyor")
    
    # Departman yetki kontrolü
    if current_user.role == 'user':
        user_depts = parse_user_departments(current_user.department)
        if department not in user_depts:
            raise HTTPException(status_code=403, detail="Bu departmana dosya yükleme yetkiniz yok")
    
    results = []
    for file in files:
        try:
            file_content = await file.read()
            
            if not file_content:
                results.append({"file": file.filename, "success": False, "error": "Dosya boş"})
                continue
            
            content, doc_type = extract_text_from_file(file.filename, file_content)
            
            if not content.strip():
                results.append({"file": file.filename, "success": False, "error": "İçerik çıkarılamadı"})
                continue
            
            success = add_document(
                content=content,
                source=file.filename,
                doc_type=doc_type,
                metadata={
                    "department": department,
                    "author": current_user.email,
                    "created_at": str(datetime.utcnow())
                }
            )
            
            results.append({
                "file": file.filename,
                "success": success,
                "type": doc_type,
                "chars": len(content) if success else 0
            })
            
        except Exception as e:
            results.append({"file": file.filename, "success": False, "error": str(e)})
    
    success_count = sum(1 for r in results if r.get('success'))
    return {
        "message": f"{success_count}/{len(files)} dosya yüklendi",
        "results": results,
        "success": success_count > 0
    }


@router.post("/documents/search", response_model=List[DocumentResponse])
async def search_docs(
    request: SearchRequest,
    current_user: User = Depends(get_current_user)
):
    """Dokümanlarda ara"""
    if not RAG_AVAILABLE:
        raise HTTPException(status_code=503, detail="RAG sistemi kullanılamıyor")
    
    results = search_documents(
        request.query, 
        n_results=request.n_results, 
        department=request.department
    )
    
    return [
        DocumentResponse(
            content=doc.get("content", ""),
            source=doc.get("source", ""),
            relevance=doc.get("relevance", 0)
        )
        for doc in results
    ]


@router.delete("/documents/{source:path}")
async def delete_doc(
    source: str,
    current_user: User = Depends(get_current_user)
):
    """Doküman sil"""
    if not RAG_AVAILABLE:
        raise HTTPException(status_code=503, detail="RAG sistemi kullanılamıyor")
    
    # Yetki kontrolü: Admin herşeyi silebilir, user sadece kendi departmanındakileri
    if current_user.role == 'user':
        # Dokümanın departmanını kontrol et
        docs = list_documents()
        doc_info = next((d for d in docs if d['source'] == source), None)
        if doc_info:
            user_depts = parse_user_departments(current_user.department)
            if doc_info['department'] not in user_depts:
                raise HTTPException(status_code=403, detail="Bu dokümanı silme yetkiniz yok")
    
    success = delete_document(source)
    
    if success:
        return {"message": f"Doküman silindi: {source}", "success": True}
    else:
        raise HTTPException(status_code=404, detail="Doküman bulunamadı")


@router.delete("/documents")
async def clear_all_docs(
    current_user: User = Depends(get_current_user)
):
    """Tüm dokümanları sil (sadece admin)"""
    if not RAG_AVAILABLE:
        raise HTTPException(status_code=503, detail="RAG sistemi kullanılamıyor")
    
    if current_user.role != "admin":
        raise HTTPException(status_code=403, detail="Sadece admin tüm dokümanları silebilir")
    
    success = clear_all_documents()
    
    if success:
        return {"message": "Tüm dokümanlar silindi", "success": True}
    else:
        raise HTTPException(status_code=500, detail="Silme işlemi başarısız")


@router.post("/teach")
async def teach_knowledge(
    request: TeachRequest,
    current_user: User = Depends(get_current_user)
):
    """Metin tabanlı bilgi öğret (Hızlı Bilgi Girişi)"""
    if not RAG_AVAILABLE:
        raise HTTPException(status_code=503, detail="RAG sistemi kullanılamıyor")
    
    # Departman yetki kontrolü
    if current_user.role == 'user':
        user_depts = parse_user_departments(current_user.department)
        if request.department not in user_depts:
            raise HTTPException(status_code=403, detail="Bu departmana bilgi ekleme yetkiniz yok")
    
    metadata = {
        "department": request.department,
        "author": current_user.email,
        "type": "manual_entry",
        "created_at": str(datetime.utcnow())
    }
    
    try:
        success = add_document(
            content=request.content,
            source=f"Manuel: {request.title}",
            doc_type="manual",
            metadata=metadata
        )
    except TypeError:
        success = add_document(
             content=request.content,
             source=f"Manuel: {request.title}",
             doc_type="text"
        )
    
    if success:
        return {"message": "Bilgi başarıyla öğretildi", "success": True}
    else:
        raise HTTPException(status_code=500, detail="Bilgi kaydedilemedi")


@router.get("/formats")
async def get_supported_formats(current_user: User = Depends(get_current_user)):
    """Desteklenen dosya formatlarını döndür"""
    return {
        "formats": list(SUPPORTED_FORMATS.keys()),
        "categories": {
            "text": [".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm", ".rtf", ".rst", ".tex", ".ini", ".cfg", ".env", ".toml", ".properties"],
            "office": [".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt", ".odt", ".ods", ".odp", ".epub"],
            "code": [".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".cs", ".cpp", ".c", ".h", ".hpp", ".sql", ".yaml", ".yml", ".go", ".rb", ".php", ".swift", ".kt", ".scala", ".rs", ".r", ".R", ".sh", ".bat", ".ps1", ".dockerfile", ".vue", ".svelte", ".graphql", ".gql", ".proto"],
            "email": [".eml", ".msg"],
            "image": [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".webp"],
            "log": [".log"]
        },
        "total_count": len(SUPPORTED_FORMATS)
    }


# ═══════════════════════════════════════════════════════════════
# URL / WEB SAYFASI ÖĞRENME
# ═══════════════════════════════════════════════════════════════

def _extract_youtube_video_id(url: str) -> Optional[str]:
    """YouTube URL'sinden video ID çıkar"""
    patterns = [
        r'(?:youtube\.com/watch\?v=|youtu\.be/|youtube\.com/embed/|youtube\.com/v/)([a-zA-Z0-9_-]{11})',
        r'youtube\.com/shorts/([a-zA-Z0-9_-]{11})',
    ]
    for pattern in patterns:
        match = re.search(pattern, url)
        if match:
            return match.group(1)
    return None


def _is_video_url(url: str) -> bool:
    """URL'nin bir video platformuna ait olup olmadığını kontrol et"""
    video_domains = ['youtube.com', 'youtu.be', 'dailymotion.com', 'vimeo.com']
    return any(domain in url.lower() for domain in video_domains)


@router.post("/learn-url")
async def learn_from_url(
    request: LearnFromUrlRequest,
    current_user: User = Depends(get_current_user)
):
    """Web sayfasından içerik öğren - URL scraping"""
    if not RAG_AVAILABLE:
        raise HTTPException(status_code=503, detail="RAG sistemi kullanılamıyor")
    
    if not WEB_SCRAPING_AVAILABLE:
        raise HTTPException(status_code=503, detail="Web scraping kütüphaneleri yüklü değil (httpx, beautifulsoup4)")
    
    # Departman yetki kontrolü
    if current_user.role == 'user':
        user_depts = parse_user_departments(current_user.department)
        if request.department not in user_depts:
            raise HTTPException(status_code=403, detail="Bu departmana içerik ekleme yetkiniz yok")
    
    # Video URL'si ise yönlendir
    if _is_video_url(request.url):
        raise HTTPException(
            status_code=400,
            detail="Bu bir video URL'si. Lütfen 'Video ile Öğren' özelliğini kullanın."
        )
    
    try:
        # Web sayfasını çek
        async with httpx.AsyncClient(timeout=30.0, follow_redirects=True) as client:
            response = await client.get(
                request.url,
                headers={
                    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
                    "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
                    "Accept-Language": "tr-TR,tr;q=0.9,en;q=0.8"
                }
            )
            response.raise_for_status()
        
        # HTML içeriğini parse et
        soup = BeautifulSoup(response.text, 'html.parser')
        
        # Gereksiz elementleri kaldır
        for tag in soup(['script', 'style', 'nav', 'footer', 'header', 'aside', 
                         'noscript', 'iframe', 'svg', 'meta', 'link']):
            tag.decompose()
        
        # Sayfa başlığını al
        page_title = request.title or (soup.title.string if soup.title else request.url)
        page_title = page_title.strip()[:200]
        
        # Ana içeriği çıkar (article, main, body sıralamasıyla)
        main_content = (
            soup.find('article') or 
            soup.find('main') or 
            soup.find('div', {'role': 'main'}) or 
            soup.find('div', class_=re.compile(r'content|article|post|entry')) or
            soup.body
        )
        
        if main_content:
            text = main_content.get_text(separator='\n', strip=True)
        else:
            text = soup.get_text(separator='\n', strip=True)
        
        # Boş satırları temizle
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        content = '\n'.join(lines)
        
        if not content or len(content) < 50:
            raise HTTPException(status_code=400, detail="Sayfadan yeterli içerik çıkarılamadı")
        
        # Çok uzun içerikleri sınırla (100K karakter)
        if len(content) > 100000:
            content = content[:100000]
            logger.info("url_content_truncated", url=request.url, original_len=len(text), truncated_to=100000)
        
        # Dokümanı ekle
        source_name = f"URL: {page_title}"
        success = add_document(
            content=content,
            source=source_name,
            doc_type="web_page",
            metadata={
                "department": request.department,
                "author": current_user.email,
                "url": request.url,
                "type": "web_page",
                "created_at": str(datetime.utcnow())
            }
        )
        
        if success:
            logger.info("url_learned", url=request.url, title=page_title, chars=len(content), user=current_user.email)
            return {
                "message": f"Web sayfası öğrenildi: {page_title}",
                "success": True,
                "title": page_title,
                "url": request.url,
                "chars": len(content),
                "type": "web_page"
            }
        else:
            raise HTTPException(status_code=500, detail="İçerik kaydedilemedi")
            
    except httpx.HTTPStatusError as e:
        raise HTTPException(status_code=400, detail=f"Sayfa erişim hatası (HTTP {e.response.status_code}): {request.url}")
    except httpx.RequestError as e:
        raise HTTPException(status_code=400, detail=f"Bağlantı hatası: {str(e)}")
    except HTTPException:
        raise
    except Exception as e:
        logger.error("url_learn_error", url=request.url, error=str(e))
        raise HTTPException(status_code=500, detail=f"URL işleme hatası: {str(e)}")


# ═══════════════════════════════════════════════════════════════
# YOUTUBE / VİDEO PLATFORMU ÖĞRENME
# ═══════════════════════════════════════════════════════════════

@router.post("/learn-video")
async def learn_from_video(
    request: LearnFromVideoRequest,
    current_user: User = Depends(get_current_user)
):
    """YouTube video altyazısından öğren"""
    if not RAG_AVAILABLE:
        raise HTTPException(status_code=503, detail="RAG sistemi kullanılamıyor")
    
    # Departman yetki kontrolü
    if current_user.role == 'user':
        user_depts = parse_user_departments(current_user.department)
        if request.department not in user_depts:
            raise HTTPException(status_code=403, detail="Bu departmana içerik ekleme yetkiniz yok")
    
    # YouTube video ID'sini çıkar
    video_id = _extract_youtube_video_id(request.url)
    if not video_id:
        raise HTTPException(
            status_code=400,
            detail="Geçerli bir YouTube URL'si giriniz. "
                   "Desteklenen formatlar: youtube.com/watch?v=XXX, youtu.be/XXX"
        )
    
    if not YOUTUBE_AVAILABLE:
        raise HTTPException(
            status_code=503, 
            detail="YouTube transcript kütüphanesi yüklü değil. "
                   "Sunucuda 'pip install youtube-transcript-api' çalıştırılmalı."
        )
    
    try:
        # Altyazı dil tercihlerini ayarla
        lang_prefs = [request.language]
        if request.language != 'en':
            lang_prefs.append('en')
        if request.language != 'tr':
            lang_prefs.append('tr')
        
        # Önce tercih edilen dilde altyazı dene
        transcript = None
        used_language = None
        
        try:
            transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
            
            # Tercih edilen dilde manuel altyazı
            for lang in lang_prefs:
                try:
                    transcript = transcript_list.find_transcript([lang])
                    used_language = lang
                    break
                except Exception:
                    continue
            
            # Manuel bulunamazsa, otomatik oluşturulmuş altyazıları dene
            if transcript is None:
                try:
                    transcript = transcript_list.find_generated_transcript(lang_prefs)
                    used_language = "auto"
                except Exception:
                    pass
            
            # Hiçbiri bulunamazsa, mevcut herhangi bir altyazıyı al
            if transcript is None:
                for t in transcript_list:
                    transcript = t
                    used_language = t.language_code
                    break
                    
        except Exception:
            # list_transcripts başarısız olursa direkt fetch dene
            try:
                raw = YouTubeTranscriptApi.get_transcript(video_id, languages=lang_prefs)
                if raw:
                    content_parts = []
                    for entry in raw:
                        text = entry.get('text', '').strip()
                        if text:
                            content_parts.append(text)
                    content = '\n'.join(content_parts)
                    used_language = request.language
                    
                    if not content or len(content) < 20:
                        raise HTTPException(status_code=400, detail="Video altyazısından yeterli içerik çıkarılamadı")
                    
                    # Dokümanı ekle
                    video_title = request.title or f"YouTube Video ({video_id})"
                    source_name = f"Video: {video_title}"
                    
                    success = add_document(
                        content=content,
                        source=source_name,
                        doc_type="video_transcript",
                        metadata={
                            "department": request.department,
                            "author": current_user.email,
                            "url": request.url,
                            "video_id": video_id,
                            "language": used_language,
                            "type": "video_transcript",
                            "created_at": str(datetime.utcnow())
                        }
                    )
                    
                    if success:
                        return {
                            "message": f"Video içeriği öğrenildi: {video_title}",
                            "success": True,
                            "title": video_title,
                            "video_id": video_id,
                            "url": request.url,
                            "chars": len(content),
                            "language": used_language,
                            "type": "video_transcript"
                        }
                    else:
                        raise HTTPException(status_code=500, detail="İçerik kaydedilemedi")
                        
            except HTTPException:
                raise
            except Exception:
                raise HTTPException(
                    status_code=400,
                    detail="Bu videonun altyazısı mevcut değil veya erişilemiyor. "
                           "Lütfen altyazısı olan bir video deneyin."
                )
        
        if transcript is None:
            raise HTTPException(
                status_code=400,
                detail="Bu videonun altyazısı bulunamadı. "
                       "Lütfen altyazısı olan bir video deneyin."
            )
        
        # Altyazı metnini al
        transcript_data = transcript.fetch()
        content_parts = []
        for entry in transcript_data:
            text = entry.get('text', '').strip()
            if text:
                content_parts.append(text)
        
        content = '\n'.join(content_parts)
        
        if not content or len(content) < 20:
            raise HTTPException(status_code=400, detail="Video altyazısından yeterli içerik çıkarılamadı")
        
        # Video başlığını al (mümkünse)
        video_title = request.title
        if not video_title:
            # Basit başlık çekme (httpx ile)
            try:
                if WEB_SCRAPING_AVAILABLE:
                    async with httpx.AsyncClient(timeout=10.0) as client:
                        resp = await client.get(
                            f"https://www.youtube.com/watch?v={video_id}",
                            headers={"User-Agent": "Mozilla/5.0"}
                        )
                        if resp.status_code == 200:
                            title_match = re.search(r'<title>(.*?)</title>', resp.text)
                            if title_match:
                                video_title = title_match.group(1).replace(' - YouTube', '').strip()
            except Exception:
                pass
            if not video_title:
                video_title = f"YouTube Video ({video_id})"
        
        # Dokümanı ekle
        source_name = f"Video: {video_title}"
        success = add_document(
            content=content,
            source=source_name,
            doc_type="video_transcript",
            metadata={
                "department": request.department,
                "author": current_user.email,
                "url": request.url,
                "video_id": video_id,
                "language": used_language,
                "type": "video_transcript",
                "created_at": str(datetime.utcnow())
            }
        )
        
        if success:
            logger.info("video_learned", video_id=video_id, title=video_title, 
                       chars=len(content), lang=used_language, user=current_user.email)
            return {
                "message": f"Video içeriği öğrenildi: {video_title}",
                "success": True,
                "title": video_title,
                "video_id": video_id,
                "url": request.url,
                "chars": len(content),
                "language": used_language,
                "type": "video_transcript"
            }
        else:
            raise HTTPException(status_code=500, detail="İçerik kaydedilemedi")
    
    except HTTPException:
        raise
    except Exception as e:
        logger.error("video_learn_error", video_id=video_id, error=str(e))
        raise HTTPException(status_code=500, detail=f"Video işleme hatası: {str(e)}")


@router.get("/capabilities")
async def get_capabilities(current_user: User = Depends(get_current_user)):
    """Sistemin mevcut öğrenme yeteneklerini döndür"""
    return {
        "file_upload": RAG_AVAILABLE,
        "web_scraping": WEB_SCRAPING_AVAILABLE,
        "youtube_transcript": YOUTUBE_AVAILABLE,
        "ocr": _check_ocr_available(),
        "supported_formats_count": len(SUPPORTED_FORMATS),
    }


def _check_ocr_available() -> bool:
    """OCR (pytesseract) mevcut mu kontrol et"""
    try:
        import pytesseract
        return True
    except ImportError:
        return False
