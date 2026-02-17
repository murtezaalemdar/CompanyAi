"""Dosya Export Servisi — Excel, PDF, PowerPoint çıktı üretimi

AI asistanın ürettiği rapor/tablo/sunum içeriklerini
kullanıcının istediği formatta dosya olarak hazırlar.
"""

import os
import re
import uuid
import time
import tempfile
from typing import Optional, Dict, List
from pathlib import Path
import structlog

logger = structlog.get_logger()

# Export dosyalarının saklanacağı dizin
EXPORT_DIR = Path(tempfile.gettempdir()) / "companyai_exports"
EXPORT_DIR.mkdir(parents=True, exist_ok=True)

# Dosya yaşam süresi (saniye) — 1 saat sonra temizlenebilir
EXPORT_TTL = 3600

# Üretilen dosya kayıtları {file_id: {path, filename, format, created_at}}
_export_registry: Dict[str, Dict] = {}


def _cleanup_old_exports():
    """1 saatten eski export dosyalarını temizle"""
    now = time.time()
    expired = [fid for fid, info in _export_registry.items()
               if now - info["created_at"] > EXPORT_TTL]
    for fid in expired:
        try:
            path = _export_registry[fid]["path"]
            if os.path.exists(path):
                os.remove(path)
            del _export_registry[fid]
        except Exception:
            pass


def detect_export_request(question: str) -> Optional[str]:
    """Kullanıcının istediği export formatını tespit eder.
    
    Returns:
        'excel', 'pdf', 'pptx', 'word', 'csv' veya None
    """
    q = question.lower()
    
    # Excel
    if any(kw in q for kw in [
        "excel", "xlsx", "xls", "tablo olarak",
        "excel dosyası", "excel formatında", "excele aktar",
        "excel olarak", "excel çıktısı", "excel raporu",
        "spreadsheet",
    ]):
        return "excel"
    
    # PDF
    if any(kw in q for kw in [
        "pdf", "pdf dosyası", "pdf formatında", "pdf olarak",
        "pdf çıktısı", "pdf raporu",
    ]):
        return "pdf"
    
    # PowerPoint
    if any(kw in q for kw in [
        "powerpoint", "pptx", "ppt", "sunum", "slayt",
        "sunum dosyası", "sunum formatında", "sunum olarak",
        "powerpoint olarak", "sunum hazırla", "slayt hazırla",
        "presentation",
    ]):
        return "pptx"
    
    # Word
    if any(kw in q for kw in [
        "word", "docx", "word dosyası", "word formatında",
        "word olarak", "word çıktısı", "word raporu",
    ]):
        return "word"
    
    # CSV
    if any(kw in q for kw in [
        "csv", "csv dosyası", "csv olarak",
    ]):
        return "csv"
    
    # Genel indirme isteği — format belirsiz, varsayılan excel
    if any(kw in q for kw in [
        "dosya olarak indir", "indirebileceğim", "export et",
        "dışa aktar", "dosya olarak ver", "dosya olarak hazırla",
        "indirmek istiyorum", "dosya halinde",
    ]):
        return "excel"
    
    return None


def _parse_content_sections(content: str) -> Dict:
    """AI yanıtını başlık, paragraf ve tablo bölümlerine ayırır."""
    sections = []
    current_section = {"title": "", "lines": []}
    
    for line in content.split("\n"):
        stripped = line.strip()
        # Markdown başlığı
        if stripped.startswith("#"):
            if current_section["lines"] or current_section["title"]:
                sections.append(current_section)
            title = stripped.lstrip("#").strip()
            current_section = {"title": title, "lines": []}
        elif stripped.startswith("**") and stripped.endswith("**"):
            if current_section["lines"] or current_section["title"]:
                sections.append(current_section)
            title = stripped.strip("*").strip()
            current_section = {"title": title, "lines": []}
        else:
            if stripped:
                current_section["lines"].append(stripped)
    
    if current_section["lines"] or current_section["title"]:
        sections.append(current_section)
    
    return sections


def _extract_tables(content: str) -> List[List[List[str]]]:
    """Markdown tablo formatını parse eder. Her tablo: [[header], [row1], [row2], ...]"""
    tables = []
    current_table = []
    
    for line in content.split("\n"):
        stripped = line.strip()
        if "|" in stripped and stripped.startswith("|"):
            cells = [c.strip() for c in stripped.split("|")[1:-1]]
            # Ayırıcı satırı atla (---|---)
            if all(re.match(r'^[-:]+$', c) for c in cells):
                continue
            current_table.append(cells)
        else:
            if current_table:
                tables.append(current_table)
                current_table = []
    
    if current_table:
        tables.append(current_table)
    
    return tables


def _clean_markdown(text: str) -> str:
    """Markdown işaretlerini temizler."""
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    text = re.sub(r'`(.*?)`', r'\1', text)
    text = re.sub(r'^#+\s*', '', text, flags=re.MULTILINE)
    text = re.sub(r'^\s*[-*]\s+', '- ', text, flags=re.MULTILINE)
    return text


# ──────────────────────────────────────────────
# Excel Export
# ──────────────────────────────────────────────

def generate_excel(content: str, title: str = "Rapor") -> Dict:
    """AI yanıtından Excel dosyası oluşturur."""
    import openpyxl
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
    
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = title[:31]  # Excel max 31 char sheet name
    
    # Stiller
    title_font = Font(name="Calibri", size=14, bold=True, color="1F4E79")
    header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
    cell_font = Font(name="Calibri", size=11)
    border = Border(
        left=Side(style='thin', color='D9E2F3'),
        right=Side(style='thin', color='D9E2F3'),
        top=Side(style='thin', color='D9E2F3'),
        bottom=Side(style='thin', color='D9E2F3'),
    )
    alt_fill = PatternFill(start_color="F2F7FB", end_color="F2F7FB", fill_type="solid")
    wrap = Alignment(wrap_text=True, vertical="top")
    
    row_idx = 1
    
    # Başlık
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=5)
    cell = ws.cell(row=1, column=1, value=title)
    cell.font = title_font
    cell.alignment = Alignment(horizontal="center", vertical="center")
    row_idx = 3
    
    # Tablolar varsa önce onları yaz
    tables = _extract_tables(content)
    if tables:
        for table in tables:
            if not table:
                continue
            # İlk satır header
            for col_idx, header in enumerate(table[0], 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=_clean_markdown(header))
                cell.font = header_font
                cell.fill = header_fill
                cell.border = border
                cell.alignment = Alignment(horizontal="center", vertical="center")
            row_idx += 1
            
            for data_row_idx, data_row in enumerate(table[1:]):
                for col_idx, val in enumerate(data_row, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=_clean_markdown(val))
                    cell.font = cell_font
                    cell.border = border
                    cell.alignment = wrap
                    if data_row_idx % 2 == 1:
                        cell.fill = alt_fill
                row_idx += 1
            row_idx += 1
    
    # Tablo yoksa, içeriği satır satır yaz
    if not tables:
        sections = _parse_content_sections(content)
        for section in sections:
            if section["title"]:
                cell = ws.cell(row=row_idx, column=1, value=section["title"])
                cell.font = Font(name="Calibri", size=12, bold=True, color="2E75B6")
                row_idx += 1
            for line in section["lines"]:
                cell = ws.cell(row=row_idx, column=1, value=_clean_markdown(line))
                cell.font = cell_font
                cell.alignment = wrap
                row_idx += 1
            row_idx += 1
    
    # Sütun genişliklerini otomatik ayarla
    for col_cells in ws.columns:
        max_len = 0
        col_letter = None
        for cell in col_cells:
            try:
                if col_letter is None:
                    col_letter = cell.column_letter
                if cell.value:
                    max_len = max(max_len, len(str(cell.value)))
            except AttributeError:
                continue  # MergedCell atla
        if col_letter:
            ws.column_dimensions[col_letter].width = min(max_len + 4, 60)
    
    # Kaydet
    file_id = str(uuid.uuid4())[:12]
    filename = f"{_safe_filename(title)}.xlsx"
    filepath = str(EXPORT_DIR / f"{file_id}_{filename}")
    wb.save(filepath)
    
    _export_registry[file_id] = {
        "path": filepath,
        "filename": filename,
        "format": "excel",
        "content_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "created_at": time.time(),
    }
    
    logger.info("export_excel_created", file_id=file_id, filename=filename)
    return {"success": True, "file_id": file_id, "filename": filename, "format": "excel"}


# ──────────────────────────────────────────────
# PDF Export
# ──────────────────────────────────────────────

def generate_pdf(content: str, title: str = "Rapor") -> Dict:
    """AI yanıtından PDF dosyası oluşturur."""
    from fpdf import FPDF
    
    class TurkishPDF(FPDF):
        def header(self):
            self.set_font("Helvetica", "B", 10)
            self.set_text_color(100, 100, 100)
            self.cell(0, 8, title[:80], align="C", new_x="LMARGIN", new_y="NEXT")
            self.line(10, self.get_y(), 200, self.get_y())
            self.ln(5)
            
        def footer(self):
            self.set_y(-15)
            self.set_font("Helvetica", "I", 8)
            self.set_text_color(150, 150, 150)
            self.cell(0, 10, f"Sayfa {self.page_no()}/{{nb}}", align="C")
    
    pdf = TurkishPDF()
    pdf.alias_nb_pages()
    pdf.set_auto_page_break(auto=True, margin=20)
    pdf.add_page()
    
    # Başlık
    pdf.set_font("Helvetica", "B", 18)
    pdf.set_text_color(31, 78, 121)
    pdf.cell(0, 15, _transliterate_turkish(title), new_x="LMARGIN", new_y="NEXT")
    pdf.ln(5)
    
    # İçerik
    sections = _parse_content_sections(content)
    
    for section in sections:
        if section["title"]:
            pdf.set_font("Helvetica", "B", 13)
            pdf.set_text_color(46, 117, 182)
            pdf.cell(0, 10, _transliterate_turkish(section["title"]), new_x="LMARGIN", new_y="NEXT")
            pdf.ln(2)
        
        pdf.set_font("Helvetica", "", 11)
        pdf.set_text_color(50, 50, 50)
        for line in section["lines"]:
            clean = _clean_markdown(line)
            clean = _transliterate_turkish(clean)
            if clean.startswith("•"):
                pdf.set_x(15)
            pdf.multi_cell(0, 6, clean, new_x="LMARGIN", new_y="NEXT")
        pdf.ln(3)
    
    # Tablolar
    tables = _extract_tables(content)
    for table in tables:
        if not table:
            continue
        col_count = len(table[0])
        col_width = min(180 / col_count, 60)
        
        # Header
        pdf.set_font("Helvetica", "B", 10)
        pdf.set_fill_color(31, 78, 121)
        pdf.set_text_color(255, 255, 255)
        for header in table[0]:
            pdf.cell(col_width, 8, _transliterate_turkish(_clean_markdown(header))[:30], border=1, fill=True, align="C")
        pdf.ln()
        
        # Data rows
        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(50, 50, 50)
        for row_i, row in enumerate(table[1:]):
            if row_i % 2 == 1:
                pdf.set_fill_color(242, 247, 251)
                fill = True
            else:
                fill = False
            for val in row:
                pdf.cell(col_width, 7, _transliterate_turkish(_clean_markdown(val))[:30], border=1, fill=fill)
            pdf.ln()
        pdf.ln(5)
    
    # Kaydet
    file_id = str(uuid.uuid4())[:12]
    filename = f"{_safe_filename(title)}.pdf"
    filepath = str(EXPORT_DIR / f"{file_id}_{filename}")
    pdf.output(filepath)
    
    _export_registry[file_id] = {
        "path": filepath,
        "filename": filename,
        "format": "pdf",
        "content_type": "application/pdf",
        "created_at": time.time(),
    }
    
    logger.info("export_pdf_created", file_id=file_id, filename=filename)
    return {"success": True, "file_id": file_id, "filename": filename, "format": "pdf"}


# ──────────────────────────────────────────────
# PowerPoint Export
# ──────────────────────────────────────────────

def generate_pptx(content: str, title: str = "Sunum") -> Dict:
    """AI yanıtından PowerPoint dosyası oluşturur."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ── Kapak slaytı ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # Arka plan
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 78, 121)
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Alt başlık
    txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4.8), Inches(9), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "CompanyAI - Kurumsal AI Asistan"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(200, 220, 240)
    p2.alignment = PP_ALIGN.CENTER
    
    # ── İçerik slaytları ──
    sections = _parse_content_sections(content)
    
    for section in sections:
        if not section["lines"] and not section["title"]:
            continue
            
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        
        # Üst şerit
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), prs.slide_width, Inches(0.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(31, 78, 121)
        shape.line.fill.background()
        
        # Slayt başlığı
        if section["title"]:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12), Inches(0.6))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = section["title"]
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
        
        # İçerik
        content_top = Inches(1.2)
        txBox = slide.shapes.add_textbox(Inches(0.8), content_top, Inches(11.5), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for i, line in enumerate(section["lines"]):
            clean = _clean_markdown(line)
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            if clean.startswith("•"):
                p.text = clean
                p.level = 1
                p.font.size = Pt(16)
            else:
                p.text = clean
                p.font.size = Pt(16)
            
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.space_after = Pt(8)
    
    # ── Tablo slaytları ──
    tables = _extract_tables(content)
    for table in tables:
        if not table or len(table) < 2:
            continue
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        rows = len(table)
        cols = len(table[0])
        
        tbl = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5 * rows)).table
        
        # Header
        for j, header in enumerate(table[0]):
            cell = tbl.cell(0, j)
            cell.text = _clean_markdown(header)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(31, 78, 121)
        
        # Data
        for i, row in enumerate(table[1:], 1):
            for j, val in enumerate(row):
                if j < cols:
                    cell = tbl.cell(i, j)
                    cell.text = _clean_markdown(val)
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(11)
                    if i % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(242, 247, 251)
    
    # Kaydet
    file_id = str(uuid.uuid4())[:12]
    filename = f"{_safe_filename(title)}.pptx"
    filepath = str(EXPORT_DIR / f"{file_id}_{filename}")
    prs.save(filepath)
    
    _export_registry[file_id] = {
        "path": filepath,
        "filename": filename,
        "format": "pptx",
        "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "created_at": time.time(),
    }
    
    logger.info("export_pptx_created", file_id=file_id, filename=filename)
    return {"success": True, "file_id": file_id, "filename": filename, "format": "pptx"}


# ──────────────────────────────────────────────
# Word Export
# ──────────────────────────────────────────────

def generate_word(content: str, title: str = "Rapor") -> Dict:
    """AI yanıtından Word dosyası oluşturur."""
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    
    doc = Document()
    
    # Stil
    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = Pt(11)
    
    # Başlık
    heading = doc.add_heading(title, level=0)
    for run in heading.runs:
        run.font.color.rgb = RGBColor(31, 78, 121)
    
    doc.add_paragraph("")  # Boşluk
    
    # İçerik
    sections = _parse_content_sections(content)
    for section in sections:
        if section["title"]:
            h = doc.add_heading(section["title"], level=2)
            for run in h.runs:
                run.font.color.rgb = RGBColor(46, 117, 182)
        
        for line in section["lines"]:
            clean = _clean_markdown(line)
            if clean.startswith("•"):
                doc.add_paragraph(clean[1:].strip(), style='List Bullet')
            else:
                doc.add_paragraph(clean)
    
    # Tablolar
    tables = _extract_tables(content)
    for table_data in tables:
        if not table_data or len(table_data) < 2:
            continue
        
        rows = len(table_data)
        cols = len(table_data[0])
        table = doc.add_table(rows=rows, cols=cols, style='Table Grid')
        
        for i, row in enumerate(table_data):
            for j, val in enumerate(row):
                if j < cols:
                    cell = table.cell(i, j)
                    cell.text = _clean_markdown(val)
                    if i == 0:
                        for p in cell.paragraphs:
                            for run in p.runs:
                                run.font.bold = True
        
        doc.add_paragraph("")
    
    # Kaydet
    file_id = str(uuid.uuid4())[:12]
    filename = f"{_safe_filename(title)}.docx"
    filepath = str(EXPORT_DIR / f"{file_id}_{filename}")
    doc.save(filepath)
    
    _export_registry[file_id] = {
        "path": filepath,
        "filename": filename,
        "format": "word",
        "content_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "created_at": time.time(),
    }
    
    logger.info("export_word_created", file_id=file_id, filename=filename)
    return {"success": True, "file_id": file_id, "filename": filename, "format": "word"}


# ──────────────────────────────────────────────
# CSV Export
# ──────────────────────────────────────────────

def generate_csv(content: str, title: str = "Veri") -> Dict:
    """AI yanıtındaki tablo verilerinden CSV dosyası oluşturur."""
    import csv
    
    tables = _extract_tables(content)
    
    file_id = str(uuid.uuid4())[:12]
    filename = f"{_safe_filename(title)}.csv"
    filepath = str(EXPORT_DIR / f"{file_id}_{filename}")
    
    with open(filepath, "w", newline="", encoding="utf-8-sig") as f:
        writer = csv.writer(f)
        if tables:
            for table in tables:
                for row in table:
                    writer.writerow([_clean_markdown(c) for c in row])
                writer.writerow([])  # Tablolar arası boşluk
        else:
            # Tablo yoksa satır satır yaz
            for line in content.split("\n"):
                clean = _clean_markdown(line.strip())
                if clean:
                    writer.writerow([clean])
    
    _export_registry[file_id] = {
        "path": filepath,
        "filename": filename,
        "format": "csv",
        "content_type": "text/csv",
        "created_at": time.time(),
    }
    
    logger.info("export_csv_created", file_id=file_id, filename=filename)
    return {"success": True, "file_id": file_id, "filename": filename, "format": "csv"}


# ──────────────────────────────────────────────
# Ana Export Fonksiyonu
# ──────────────────────────────────────────────

def generate_export(content: str, fmt: str, title: str = "Rapor") -> Optional[Dict]:
    """İstenen formatta export üretir.
    
    Args:
        content: AI yanıt metni
        fmt: 'excel', 'pdf', 'pptx', 'word', 'csv'
        title: Dosya/rapor başlığı
    
    Returns:
        {file_id, filename, format} veya hata durumunda None
    """
    _cleanup_old_exports()
    
    generators = {
        "excel": generate_excel,
        "pdf": generate_pdf,
        "pptx": generate_pptx,
        "word": generate_word,
        "csv": generate_csv,
    }
    
    gen = generators.get(fmt)
    if not gen:
        logger.error("export_unknown_format", format=fmt)
        return None
    
    try:
        return gen(content, title)
    except Exception as e:
        logger.error("export_generation_error", format=fmt, error=str(e))
        return None


def get_export_info(file_id: str) -> Optional[Dict]:
    """Export dosyası bilgisini döndürür."""
    return _export_registry.get(file_id)


# ──────────────────────────────────────────────
# Yardımcı
# ──────────────────────────────────────────────

# Türkçe → ASCII transliteration (PDF font desteği olmayan karakterler için)
_TR_MAP = str.maketrans({
    'ı': 'i', 'İ': 'I', 'ş': 's', 'Ş': 'S',
    'ç': 'c', 'Ç': 'C', 'ğ': 'g', 'Ğ': 'G',
    'ö': 'o', 'Ö': 'O', 'ü': 'u', 'Ü': 'U',
})


def _safe_filename(name: str) -> str:
    """Dosya adı için güvenli karakter seti — HTTP header latin-1 uyumlu."""
    # Türkçe karakterleri ASCII karşılıklarına çevir (header uyumluluğu)
    safe = name.translate(_TR_MAP)
    safe = re.sub(r'[^\w\s-]', '', safe)
    safe = re.sub(r'\s+', '_', safe.strip())
    return safe[:50] or "rapor"


def _transliterate_turkish(text: str) -> str:
    """Türkçe karakterleri ASCII karşılıklarına dönüştürür (PDF uyumluluğu).
    Helvetica fontunda bulunmayan tüm karakterleri temizler."""
    result = text.translate(_TR_MAP)
    # Sadece ASCII + temel Latin aralığında karakterleri bırak
    cleaned = []
    for ch in result:
        code = ord(ch)
        if code < 256:  # Latin-1 aralığı (Helvetica destekler)
            cleaned.append(ch)
        elif ch in '•–—""''…€™':
            # Yaygın unicode → ASCII karşılıkları
            replacements = {'•': '-', '–': '-', '—': '-', '\u201c': '"', '\u201d': '"',
                          '\u2018': "'", '\u2019': "'", '…': '...', '€': 'EUR', '™': '(TM)'}
            cleaned.append(replacements.get(ch, ''))
        else:
            cleaned.append(' ')  # Bilinmeyen karakteri boşluk yap
    return ''.join(cleaned)


# Format etiketleri (frontend gösterimi için)
FORMAT_LABELS = {
    "excel": {"label": "Excel", "icon": "📊", "ext": ".xlsx"},
    "pdf": {"label": "PDF", "icon": "📄", "ext": ".pdf"},
    "pptx": {"label": "PowerPoint", "icon": "📽️", "ext": ".pptx"},
    "word": {"label": "Word", "icon": "📝", "ext": ".docx"},
    "csv": {"label": "CSV", "icon": "📋", "ext": ".csv"},
}
